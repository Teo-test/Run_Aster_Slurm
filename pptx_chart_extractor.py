#!/usr/bin/env python3
"""
PPTX Chart Extractor — Extrait les données des graphes d'un PowerPoint
Exporte chaque graphe en CSV et génère un plot matplotlib.

Usage: python pptx_chart_extractor.py [fichier.pptx]
"""

import sys
import os
import re
import argparse
from pathlib import Path

# ── Dépendances ───────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    print("[ERREUR] python-pptx manquant → pip install python-pptx")
    sys.exit(1)

try:
    import pandas as pd
    import matplotlib.pyplot as plt
    import matplotlib.cm as cm
    import numpy as np
except ImportError as e:
    print(f"[ERREUR] Dépendance manquante : {e}")
    print("Installe : pip install pandas matplotlib numpy")
    sys.exit(1)

# ─── Couleurs terminal (même style que csv_comparator) ────────────────────────
class C:
    BOLD    = "\033[1m"
    DIM     = "\033[2m"
    CYAN    = "\033[96m"
    GREEN   = "\033[92m"
    YELLOW  = "\033[93m"
    RED     = "\033[91m"
    BLUE    = "\033[94m"
    MAGENTA = "\033[95m"
    RESET   = "\033[0m"

def titre(texte):
    largeur = 60
    print(f"\n{C.CYAN}{C.BOLD}{'─' * largeur}{C.RESET}")
    print(f"{C.CYAN}{C.BOLD}  {texte}{C.RESET}")
    print(f"{C.CYAN}{C.BOLD}{'─' * largeur}{C.RESET}")

def info(texte):   print(f"  {C.BLUE}ℹ {C.RESET}{texte}")
def ok(texte):     print(f"  {C.GREEN}✔ {C.RESET}{texte}")
def warn(texte):   print(f"  {C.YELLOW}⚠ {C.RESET}{texte}")
def erreur(texte): print(f"  {C.RED}✘ {C.RESET}{texte}")

def menu_numerote(titre_menu, options, allow_multiple=False):
    print(f"\n{C.BOLD}  {titre_menu}{C.RESET}")
    for i, opt in enumerate(options, 1):
        print(f"    {C.CYAN}{i:>2}{C.RESET}. {opt}")
    print()
    while True:
        if allow_multiple:
            entree = input(f"  {C.BOLD}Choix (ex: 1 3 5 ou 'all'){C.RESET} : ").strip()
            if entree.lower() == "all":
                return list(range(len(options)))
            try:
                choix = [int(x) - 1 for x in entree.split()]
                if all(0 <= c < len(options) for c in choix) and choix:
                    return choix
            except ValueError:
                pass
        else:
            entree = input(f"  {C.BOLD}Choix{C.RESET} : ").strip()
            try:
                c = int(entree) - 1
                if 0 <= c < len(options):
                    return c
            except ValueError:
                pass
        erreur("Entrée invalide, réessaie.")

def input_prompt(texte, defaut=None):
    if defaut:
        rep = input(f"  {C.BOLD}{texte}{C.RESET} [{C.DIM}{defaut}{C.RESET}] : ").strip()
        return rep if rep else defaut
    return input(f"  {C.BOLD}{texte}{C.RESET} : ").strip()

def couleurs(n):
    return cm.tab10(np.linspace(0, 0.9, max(n, 1)))

def slugify(texte):
    """Transforme un titre en nom de fichier propre."""
    texte = re.sub(r'[^\w\s-]', '', texte.lower())
    return re.sub(r'[\s-]+', '_', texte).strip('_') or "graphe"

# ─── Type de graphe ───────────────────────────────────────────────────────────
# Correspondance type XL → famille lisible
FAMILLES = {
    # Lignes
    "LINE": "ligne",
    "LINE_MARKERS": "ligne",
    "LINE_STACKED": "ligne",
    "LINE_STACKED_100": "ligne",
    "LINE_MARKERS_STACKED": "ligne",
    "LINE_MARKERS_STACKED_100": "ligne",
    # Barres / colonnes
    "BAR_CLUSTERED": "barres",
    "BAR_STACKED": "barres",
    "BAR_STACKED_100": "barres",
    "COLUMN_CLUSTERED": "barres",
    "COLUMN_STACKED": "barres",
    "COLUMN_STACKED_100": "barres",
    # Scatter / Bubble
    "XY_SCATTER": "scatter",
    "XY_SCATTER_LINES": "scatter",
    "XY_SCATTER_LINES_NO_MARKERS": "scatter",
    "XY_SCATTER_SMOOTH": "scatter",
    "XY_SCATTER_SMOOTH_NO_MARKERS": "scatter",
    "BUBBLE": "scatter",
    "BUBBLE_THREE_D_EFFECT": "scatter",
    # Camembert
    "PIE": "pie",
    "PIE_EXPLODED": "pie",
    "PIE_OF_PIE": "pie",
    "BAR_OF_PIE": "pie",
    "DOUGHNUT": "pie",
    "DOUGHNUT_EXPLODED": "pie",
    # Aires
    "AREA": "aire",
    "AREA_STACKED": "aire",
    "AREA_STACKED_100": "aire",
}

def famille_chart(chart):
    try:
        nom = chart.chart_type.name
        return FAMILLES.get(nom, "barres"), nom
    except Exception:
        return "barres", "INCONNU"

# ─── Extraction des données ───────────────────────────────────────────────────
def extraire_chart_data(chart):
    """
    Extrait les données d'un graphe pptx en DataFrame.
    Gère : catégories/séries classiques + XY scatter.
    """
    famille, type_nom = famille_chart(chart)
    rows = []
    colonnes = []

    try:
        # ── Scatter / Bubble (pas de catégories textuelles)
        if famille == "scatter":
            for serie in chart.series:
                nom_serie = serie.name or "Série"
                xs = list(serie.values)        # ordonnées X
                ys = list(serie.values)        # on va chercher Y séparément
                # python-pptx expose x_values et y_values pour scatter
                try:
                    xs = list(serie.x_values)
                    ys = list(serie.y_values)
                except AttributeError:
                    pass
                for x, y in zip(xs, ys):
                    rows.append({"série": nom_serie, "x": x, "y": y})
            colonnes = ["série", "x", "y"]
            df = pd.DataFrame(rows, columns=colonnes) if rows else pd.DataFrame()
            return df, famille, type_nom

        # ── Pie / Doughnut
        if famille == "pie":
            plot = chart.plots[0]
            serie = plot.series[0]
            try:
                cats = [str(c.label) for c in serie.data_labels] if serie.data_labels else []
            except Exception:
                cats = []
            try:
                cats = [str(p.label) for p in serie.points]
            except Exception:
                pass
            try:
                # Catégories via chart.plots[0].categories
                cats = [str(c) for c in chart.plots[0].categories]
            except Exception:
                pass
            vals = list(serie.values)
            if not cats:
                cats = [f"Catégorie {i+1}" for i in range(len(vals))]
            for cat, val in zip(cats, vals):
                rows.append({"catégorie": cat, "valeur": val})
            df = pd.DataFrame(rows)
            return df, famille, type_nom

        # ── Cas général : catégories + séries (lignes, barres, aires)
        plot = chart.plots[0]

        # Catégories (axe X)
        try:
            cats = [str(c) for c in plot.categories]
        except Exception:
            cats = None

        for serie in plot.series:
            nom_serie = serie.name or "Série"
            vals = list(serie.values)
            n = len(vals)
            if cats and len(cats) == n:
                cat_list = cats
            else:
                cat_list = [str(i) for i in range(n)]
            for cat, val in zip(cat_list, vals):
                rows.append({"catégorie": cat, "série": nom_serie, "valeur": val})

        df = pd.DataFrame(rows)
        if not df.empty:
            # Pivot : une colonne par série
            try:
                df = df.pivot_table(index="catégorie", columns="série", values="valeur", aggfunc="first").reset_index()
                df.columns.name = None
            except Exception:
                pass
        return df, famille, type_nom

    except Exception as e:
        warn(f"Extraction partielle : {e}")
        return pd.DataFrame(), famille, type_nom

# ─── Analyse du fichier PPTX ─────────────────────────────────────────────────
def analyser_pptx(chemin):
    """Parcourt toutes les slides et collecte les graphes."""
    prs = Presentation(chemin)
    graphes = []
    graphe_global_idx = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        # Titre de la slide (si disponible)
        titre_slide = ""
        try:
            titre_slide = slide.shapes.title.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            graphe_global_idx += 1

            # Titre du graphe
            titre_graphe = ""
            try:
                titre_graphe = chart.chart_title.text_frame.text.strip()
            except Exception:
                pass
            if not titre_graphe:
                titre_graphe = titre_slide or f"Graphe {graphe_global_idx}"

            famille, type_nom = famille_chart(chart)
            df, _, _ = extraire_chart_data(chart)

            graphes.append({
                "idx":           graphe_global_idx,
                "slide":         slide_idx,
                "titre":         titre_graphe,
                "famille":       famille,
                "type_xl":       type_nom,
                "df":            df,
                "n_series":      len(df.columns) - 1 if not df.empty and "série" not in df.columns else 1,
                "n_points":      len(df),
            })

    return graphes

# ─── Affichage résumé ────────────────────────────────────────────────────────
def afficher_resume(graphes):
    titre("GRAPHES TROUVÉS")
    if not graphes:
        warn("Aucun graphe détecté dans ce fichier.")
        return
    print(f"  {'N°':>3}  {'Slide':>5}  {'Type':<10}  {'Points':>6}  Titre")
    print(f"  {'─'*3}  {'─'*5}  {'─'*10}  {'─'*6}  {'─'*30}")
    for g in graphes:
        vide = f" {C.YELLOW}[vide]{C.RESET}" if g["df"].empty else ""
        print(f"  {C.CYAN}{g['idx']:>3}{C.RESET}  "
              f"  {g['slide']:>3}  "
              f"  {C.BOLD}{g['famille']:<10}{C.RESET}  "
              f"  {g['n_points']:>5}  "
              f"  {g['titre']}{vide}")

# ─── Export CSV ──────────────────────────────────────────────────────────────
def exporter_csv(graphes, selection, dossier_out):
    titre("EXPORT CSV")
    dossier = Path(dossier_out)
    dossier.mkdir(parents=True, exist_ok=True)

    for idx in selection:
        g = graphes[idx]
        if g["df"].empty:
            warn(f"Graphe {g['idx']} ({g['titre']}) : données vides, ignoré.")
            continue
        nom_fichier = f"graphe_{g['idx']:02d}_{slugify(g['titre'])}.csv"
        chemin = dossier / nom_fichier
        g["df"].to_csv(chemin, index=False, encoding="utf-8-sig")
        ok(f"Exporté : {C.BOLD}{chemin}{C.RESET}  ({len(g['df'])} lignes)")

# ─── Plots ───────────────────────────────────────────────────────────────────
def style_ax(ax, titre_graphe, xlabel="", ylabel=""):
    ax.set_title(titre_graphe, fontsize=12, fontweight='bold', pad=10)
    if xlabel: ax.set_xlabel(xlabel, fontsize=9)
    if ylabel: ax.set_ylabel(ylabel, fontsize=9)
    ax.grid(True, alpha=0.25, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

def plot_graphe(g, sauvegarder=False, dossier_out="."):
    df = g["df"]
    famille = g["famille"]
    titre_g = g["titre"]
    fig, ax = plt.subplots(figsize=(9, 5))

    if df.empty:
        ax.text(0.5, 0.5, "Données non disponibles", ha='center', va='center',
                transform=ax.transAxes, fontsize=14, color='gray')
        ax.set_title(titre_g, fontsize=12, fontweight='bold')
        plt.tight_layout()
        if sauvegarder:
            _sauvegarder_fig(fig, g, dossier_out)
        else:
            plt.show()
        return

    # ── Scatter ────────────────────────────────────────────────────────────────
    if famille == "scatter":
        series_noms = df["série"].unique() if "série" in df.columns else ["Données"]
        cols = couleurs(len(series_noms))
        for i, nom in enumerate(series_noms):
            sous = df[df["série"] == nom] if "série" in df.columns else df
            ax.scatter(sous["x"], sous["y"], label=nom, color=cols[i], alpha=0.75, s=40, edgecolors='none')
        style_ax(ax, titre_g, "X", "Y")
        if len(series_noms) > 1: ax.legend(fontsize=8)

    # ── Pie / Doughnut ─────────────────────────────────────────────────────────
    elif famille == "pie":
        if "catégorie" in df.columns and "valeur" in df.columns:
            vals = df["valeur"].dropna()
            labs = df["catégorie"][vals.index]
            ax.pie(vals, labels=labs, autopct='%1.1f%%', startangle=90,
                   colors=couleurs(len(vals)), pctdistance=0.8,
                   wedgeprops=dict(linewidth=0.5, edgecolor='white'))
            ax.set_title(titre_g, fontsize=12, fontweight='bold', pad=10)
        else:
            ax.text(0.5, 0.5, "Structure inattendue", ha='center', va='center', transform=ax.transAxes)

    # ── Barres ─────────────────────────────────────────────────────────────────
    elif famille == "barres":
        col_cat = df.columns[0]
        cols_val = df.columns[1:]
        x = np.arange(len(df))
        n = len(cols_val)
        width = 0.8 / max(n, 1)
        clrs = couleurs(n)
        for i, col in enumerate(cols_val):
            offset = (i - n / 2 + 0.5) * width
            ax.bar(x + offset, df[col], width=width * 0.92, label=col, color=clrs[i], alpha=0.85)
        ax.set_xticks(x)
        ax.set_xticklabels(df[col_cat], rotation=30, ha='right', fontsize=8)
        style_ax(ax, titre_g, col_cat, "Valeur")
        if n > 1: ax.legend(fontsize=8)

    # ── Ligne / Aire ───────────────────────────────────────────────────────────
    else:
        col_cat = df.columns[0]
        cols_val = df.columns[1:]
        clrs = couleurs(len(cols_val))
        for i, col in enumerate(cols_val):
            vals = pd.to_numeric(df[col], errors='coerce')
            if famille == "aire":
                ax.fill_between(range(len(df)), vals, alpha=0.4, color=clrs[i])
            ax.plot(range(len(df)), vals, label=col, color=clrs[i], linewidth=1.8,
                    marker='o', markersize=3)
        ax.set_xticks(range(len(df)))
        ax.set_xticklabels(df[col_cat], rotation=30, ha='right', fontsize=8)
        style_ax(ax, titre_g, col_cat, "Valeur")
        if len(cols_val) > 1: ax.legend(fontsize=8)

    plt.tight_layout()
    if sauvegarder:
        _sauvegarder_fig(fig, g, dossier_out)
    else:
        plt.show()

def _sauvegarder_fig(fig, g, dossier_out):
    dossier = Path(dossier_out)
    dossier.mkdir(parents=True, exist_ok=True)
    nom = f"graphe_{g['idx']:02d}_{slugify(g['titre'])}.png"
    chemin = dossier / nom
    fig.savefig(chemin, dpi=150, bbox_inches='tight')
    plt.close(fig)
    ok(f"Image sauvegardée : {C.BOLD}{chemin}{C.RESET}")

# ─── Menus ───────────────────────────────────────────────────────────────────
def choisir_graphes(graphes, message="Quels graphes ?"):
    labels = [f"[Slide {g['slide']}] {g['famille']:<8} — {g['titre']}" for g in graphes]
    idxs = menu_numerote(message, labels + ["Tous"], allow_multiple=True)
    if len(graphes) in idxs:  # "Tous" sélectionné
        return list(range(len(graphes)))
    return idxs

def menu_exporter(graphes, dossier_out):
    selection = choisir_graphes(graphes, "Exporter quels graphes en CSV ?")
    exporter_csv(graphes, selection, dossier_out)

def menu_plot(graphes, dossier_out):
    selection = choisir_graphes(graphes, "Afficher quels graphes ?")

    mode_opts = ["Afficher à l'écran", "Sauvegarder en PNG", "Les deux"]
    mode_idx = menu_numerote("Mode d'affichage", mode_opts)

    for idx in selection:
        g = graphes[idx]
        info(f"Graphe {g['idx']} : {g['titre']}  ({g['type_xl']})")
        if g["df"].empty:
            warn("  Données vides — plot ignoré.")
            continue
        if mode_idx == 0:
            plot_graphe(g, sauvegarder=False)
        elif mode_idx == 1:
            plot_graphe(g, sauvegarder=True, dossier_out=dossier_out)
        else:
            plot_graphe(g, sauvegarder=True, dossier_out=dossier_out)
            plot_graphe(g, sauvegarder=False)

def menu_tout_exporter(graphes, dossier_out):
    """Export CSV + PNG de tous les graphes d'un coup."""
    titre("EXPORT COMPLET (CSV + PNG)")
    non_vides = [i for i, g in enumerate(graphes) if not g["df"].empty]
    if not non_vides:
        warn("Aucun graphe avec des données à exporter.")
        return
    info(f"{len(non_vides)} graphe(s) avec données → export vers {C.BOLD}{dossier_out}/{C.RESET}")
    exporter_csv(graphes, non_vides, dossier_out)
    for idx in non_vides:
        plot_graphe(graphes[idx], sauvegarder=True, dossier_out=dossier_out)
    ok(f"Export terminé dans {C.BOLD}{dossier_out}/{C.RESET}")

def afficher_detail(graphes):
    titre("DÉTAIL D'UN GRAPHE")
    idx = menu_numerote("Quel graphe ?", [f"[Slide {g['slide']}] {g['titre']}" for g in graphes])
    g = graphes[idx]
    print(f"\n  {C.BOLD}{C.MAGENTA}{g['titre']}{C.RESET}")
    print(f"  Slide       : {g['slide']}")
    print(f"  Type XL     : {g['type_xl']}")
    print(f"  Famille     : {g['famille']}")
    print(f"  Lignes      : {len(g['df'])}")
    if not g["df"].empty:
        print(f"  Colonnes    : {', '.join(g['df'].columns)}")
        print(f"\n{g['df'].head(8).to_string(index=False)}\n")
    else:
        warn("  Aucune donnée extractible.")

# ─── Menu principal ───────────────────────────────────────────────────────────
MENU_PRINCIPAL = [
    "Afficher / sauvegarder des graphes",
    "Exporter des graphes en CSV",
    "Export complet (tous CSV + PNG)",
    "Voir le détail d'un graphe",
    "Résumé des graphes",
    "Charger un autre fichier PPTX",
    "Quitter",
]

def charger_pptx(chemin):
    p = Path(chemin)
    if not p.exists():
        erreur(f"Fichier introuvable : {chemin}")
        return None, None
    info(f"Analyse de {C.BOLD}{p.name}{C.RESET} …")
    graphes = analyser_pptx(chemin)
    ok(f"{len(graphes)} graphe(s) trouvé(s) dans {p.name}")
    return graphes, p

def main():
    parser = argparse.ArgumentParser(description="Extracteur de graphes PPTX → CSV + plots")
    parser.add_argument("fichier", nargs="?", help="Fichier .pptx à analyser")
    parser.add_argument("-o", "--output", default="pptx_export", help="Dossier de sortie (défaut: pptx_export)")
    args = parser.parse_args()

    print(f"\n{C.CYAN}{C.BOLD}")
    print("  ╔══════════════════════════════════════════╗")
    print("  ║      PPTX CHART EXTRACTOR  v1.0          ║")
    print("  ║   Extraction de graphes PowerPoint       ║")
    print("  ╚══════════════════════════════════════════╝")
    print(C.RESET)

    graphes = []
    chemin_pptx = None

    if args.fichier:
        graphes, p = charger_pptx(args.fichier)
        if graphes is None:
            sys.exit(1)
        chemin_pptx = args.fichier
        afficher_resume(graphes)
    else:
        info("Aucun fichier spécifié. Utilise le menu pour en charger un.")
        info("Astuce : python pptx_chart_extractor.py presentation.pptx")

    dossier_out = args.output

    while True:
        titre("MENU PRINCIPAL")
        if chemin_pptx:
            info(f"Fichier actif : {C.BOLD}{Path(chemin_pptx).name}{C.RESET}  "
                 f"({len(graphes)} graphe(s))  →  sortie : {C.BOLD}{dossier_out}/{C.RESET}")
        else:
            warn("Aucun fichier chargé")

        choix = menu_numerote("Que veux-tu faire ?", MENU_PRINCIPAL)

        if not graphes and choix not in (5, 6):
            warn("Charge d'abord un fichier PPTX (option 6).")
            continue

        if choix == 0:
            menu_plot(graphes, dossier_out)
        elif choix == 1:
            menu_exporter(graphes, dossier_out)
        elif choix == 2:
            menu_tout_exporter(graphes, dossier_out)
        elif choix == 3:
            afficher_detail(graphes)
        elif choix == 4:
            afficher_resume(graphes)
        elif choix == 5:
            chemin = input_prompt("Chemin du fichier PPTX").strip('"').strip("'")
            graphes, p = charger_pptx(chemin)
            if graphes is not None:
                chemin_pptx = chemin
                afficher_resume(graphes)
        elif choix == 6:
            print(f"\n  {C.GREEN}Au revoir !{C.RESET}\n")
            break

if __name__ == "__main__":
    main()
